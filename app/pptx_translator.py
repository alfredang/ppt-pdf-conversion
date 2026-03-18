"""PPTX Chinese-to-English translation module.

Translates text in-place within an existing PPTX file,
preserving all formatting, layouts, images, and animations.
"""

from pptx import Presentation
from pptx.util import Pt
from app.translator import Translator, has_cjk


def translate_pptx(input_path, output_path, progress_callback=None):
    """Translate Chinese text in a PPTX to English, keeping formatting intact.

    Args:
        input_path: Path to input PPTX file
        output_path: Path for output PPTX file
        progress_callback: Optional callable(message: str) for progress updates
    """
    def _progress(msg):
        if progress_callback:
            progress_callback(msg)

    _progress("Opening presentation...")
    prs = Presentation(input_path)
    translator = Translator()
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        _progress(f"Translating slide {i + 1}/{total_slides}")
        _translate_slide(slide, translator)

    _progress("Saving translated presentation...")
    prs.save(output_path)
    _progress("Conversion complete!")
    return output_path


def _translate_slide(slide, translator):
    """Translate all text in a slide's shapes."""
    for shape in slide.shapes:
        _translate_shape(shape, translator)

    # Also handle notes slide if present
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            _translate_text_frame(notes_frame, translator)


def _translate_shape(shape, translator):
    """Translate text within a shape, handling groups and tables."""
    # Group shapes contain child shapes
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            _translate_shape(child, translator)
        return

    # Tables have cells with text frames
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _translate_text_frame(cell.text_frame, translator)
        return

    # Regular shapes with text frames
    if shape.has_text_frame:
        _translate_text_frame(shape.text_frame, translator)


def _translate_text_frame(text_frame, translator):
    """Translate text in a text frame, preserving per-run formatting."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text or not run.text.strip():
                continue
            if has_cjk(run.text):
                translated = translator.translate_text(run.text)
                if translated:
                    run.text = translated
