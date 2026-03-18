"""PPTX generation module using python-pptx."""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Conversion: PDF points to EMU (1 point = 12700 EMU)
PT_TO_EMU = 12700

# Common CJK font to Latin font mapping
CJK_FONT_MAP = {
    'SimSun': 'Times New Roman',
    'SimHei': 'Arial',
    'NSimSun': 'Times New Roman',
    'FangSong': 'Times New Roman',
    'KaiTi': 'Georgia',
    'Microsoft YaHei': 'Segoe UI',
    'Microsoft JhengHei': 'Segoe UI',
    'DengXian': 'Segoe UI',
    'STSong': 'Times New Roman',
    'STHeiti': 'Arial',
    'STKaiti': 'Georgia',
    'STFangsong': 'Times New Roman',
    'PingFang SC': 'Helvetica Neue',
    'PingFang TC': 'Helvetica Neue',
    'Heiti SC': 'Arial',
    'Songti SC': 'Times New Roman',
    'Hiragino Sans GB': 'Helvetica Neue',
}


def _map_font(pdf_font_name):
    """Map a PDF font name to an appropriate Latin font.

    Preserves the original font if it's already a Latin font.
    Maps known CJK fonts to visually similar Latin alternatives.
    """
    if not pdf_font_name:
        return 'Arial'

    # Strip common PDF font prefixes like "ABCDEF+"
    clean = pdf_font_name
    if '+' in clean:
        clean = clean.split('+', 1)[1]

    # Check CJK mapping (try exact match, then partial match)
    for cjk_name, latin_name in CJK_FONT_MAP.items():
        if cjk_name.lower() in clean.lower():
            return latin_name

    # Strip style suffixes to get base font name
    base = clean
    for suffix in ('-Bold', '-Italic', '-BoldItalic', ',Bold', ',Italic',
                    '-Regular', ',Regular', '-Light', '-Medium', '-Semibold'):
        base = base.replace(suffix, '')

    return base


def _color_from_int(color_int):
    """Convert a PDF integer color to pptx RGBColor."""
    if color_int is None or color_int == 0:
        return RGBColor(0, 0, 0)  # Default black
    r = (color_int >> 16) & 0xFF
    g = (color_int >> 8) & 0xFF
    b = color_int & 0xFF
    return RGBColor(r, g, b)


def create_presentation(pages_data, progress_callback=None):
    """Create a PPTX presentation from extracted and translated page data."""
    prs = Presentation()

    # Use first page dimensions to set slide size
    if pages_data:
        pw = pages_data[0].get('width', 960)
        ph = pages_data[0].get('height', 540)
        prs.slide_width = Emu(int(pw * PT_TO_EMU))
        prs.slide_height = Emu(int(ph * PT_TO_EMU))

    for i, page_data in enumerate(pages_data):
        if progress_callback:
            progress_callback(f"Generating slide {i + 1}/{len(pages_data)}")

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        _add_images(slide, page_data)
        _add_text_groups(slide, page_data, prs.slide_width)

    return prs


def _add_images(slide, page_data):
    """Add images to a slide."""
    for img in page_data.get('images', []):
        bbox = img['bbox']
        left = Emu(int(bbox[0] * PT_TO_EMU))
        top = Emu(int(bbox[1] * PT_TO_EMU))
        width = Emu(int((bbox[2] - bbox[0]) * PT_TO_EMU))
        height = Emu(int((bbox[3] - bbox[1]) * PT_TO_EMU))

        stream = img['stream']
        stream.seek(0)

        try:
            slide.shapes.add_picture(stream, left, top, width, height)
        except Exception as e:
            print(f"Error adding image: {e}")


def _is_centered(bbox, page_width, tolerance=0.15):
    """Check if a text box is roughly centered on the page."""
    text_center = (bbox[0] + bbox[2]) / 2
    page_center = page_width / 2
    return abs(text_center - page_center) / page_width < tolerance


def _add_text_groups(slide, page_data, slide_width):
    """Add text groups (lines) as text boxes on the slide."""
    page_width = page_data.get('width', 960)
    groups = page_data.get('text_groups', [])

    # Find the max font size to help identify titles
    max_font_size = max((g['font_size'] for g in groups), default=12)

    for group in groups:
        text = group.get('translated_text', group['text'])
        bbox = group['bbox']
        font_size = group['font_size']

        # Detect if this is a title-like element (large font + centered in original)
        is_title = (font_size >= max_font_size * 0.9
                    and _is_centered(bbox, page_width))

        if is_title:
            # Center title across full slide width with padding
            padding = Emu(int(page_width * 0.05 * PT_TO_EMU))
            left = padding
            top = Emu(int(bbox[1] * PT_TO_EMU))
            width = slide_width - 2 * padding
            height = Emu(int((bbox[3] - bbox[1]) * PT_TO_EMU))
        else:
            left = Emu(int(bbox[0] * PT_TO_EMU))
            top = Emu(int(bbox[1] * PT_TO_EMU))
            width = Emu(int((bbox[2] - bbox[0]) * PT_TO_EMU))
            height = Emu(int((bbox[3] - bbox[1]) * PT_TO_EMU))

        # Ensure minimum dimensions
        if width < Emu(50000):
            width = Emu(500000)
        if height < Emu(50000):
            height = Emu(300000)

        try:
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = is_title

            p = tf.paragraphs[0]
            if is_title:
                p.alignment = PP_ALIGN.CENTER

            run = p.add_run()
            run.text = text

            # Apply formatting
            font = run.font
            font.size = Pt(font_size)
            font.color.rgb = _color_from_int(group['color'])
            font.bold = group.get('is_bold', False)
            font.italic = group.get('is_italic', False)
            font.name = _map_font(group.get('font_name'))

            # Remove default margins for tighter positioning
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)

        except Exception as e:
            print(f"Error adding text '{text[:30]}...': {e}")
